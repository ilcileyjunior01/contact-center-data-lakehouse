"""
Gerador de apresentação PowerPoint
Contact Center Data Lakehouse — AWS
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
import copy

# ── Paleta de cores ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # Azul muito escuro
BG_CARD      = RGBColor(0x14, 0x2A, 0x42)   # Azul médio (cards)
ACCENT_BLUE  = RGBColor(0x00, 0x9B, 0xFF)   # Azul elétrico
ACCENT_TEAL  = RGBColor(0x00, 0xD4, 0xAA)   # Verde-água
ACCENT_ORG   = RGBColor(0xFF, 0x99, 0x00)   # Laranja AWS
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)
TEXT_MUTED   = RGBColor(0x7F, 0x9E, 0xBB)
OK_GREEN     = RGBColor(0x2E, 0xCC, 0x71)
WARN_YELLOW  = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.33)   # largura widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=BG_CARD, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width_pt=2):
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector


def accent_bar(slide, x=0.4, y=1.05, w=0.08, h=0.45, color=ACCENT_BLUE):
    box(slide, x, y, w, h, fill_color=color)


def header(slide, title, subtitle=None, y_title=0.28, y_sub=0.72):
    txt(slide, title, 0.5, y_title, 12.3, 0.6,
        size=32, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.5, y_sub, 12.3, 0.4,
            size=16, color=TEXT_MUTED, align=PP_ALIGN.LEFT)
    # linha separadora
    box(slide, 0.5, 1.18, 12.33, 0.04, fill_color=ACCENT_BLUE)


def pill(slide, text, x, y, w=1.8, h=0.38,
         bg_color=ACCENT_BLUE, font_size=13):
    b = box(slide, x, y, w, h, fill_color=bg_color)
    txt(slide, text, x + 0.05, y + 0.04, w - 0.1, h - 0.08,
        size=font_size, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    return b


def card(slide, title, body_lines, x, y, w=3.8, h=2.2,
         title_color=ACCENT_BLUE):
    box(slide, x, y, w, h, fill_color=BG_CARD)
    txt(slide, title, x + 0.15, y + 0.12, w - 0.3, 0.4,
        size=14, bold=True, color=title_color)
    box(slide, x + 0.15, y + 0.52, w - 0.3, 0.03, fill_color=title_color)
    body = "\n".join(body_lines)
    txt(slide, body, x + 0.15, y + 0.62, w - 0.3, h - 0.75,
        size=12, color=TEXT_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Capa"""
    sl = blank_slide(prs)
    bg(sl, BG_DARK)

    # Fundo decorativo — bloco lateral esquerdo
    box(sl, 0, 0, 4.2, 7.5, fill_color=RGBColor(0x09, 0x12, 0x1E))

    # Linha vertical accent
    box(sl, 4.0, 0, 0.07, 7.5, fill_color=ACCENT_BLUE)

    # Badges de tecnologia (lado esquerdo)
    badges = [
        ("AWS",         ACCENT_ORG),
        ("PySpark",     RGBColor(0xE2, 0x5A, 0x1E)),
        ("Airflow",     RGBColor(0x01, 0x7C, 0xEE)),
        ("Iceberg v2",  RGBColor(0x4A, 0x90, 0xD9)),
        ("Redshift",    RGBColor(0x8C, 0x4F, 0xFF)),
        ("Glue 4.0",    ACCENT_ORG),
    ]
    for i, (label, color) in enumerate(badges):
        yb = 1.5 + i * 0.72
        pill(sl, label, 0.4, yb, w=3.0, h=0.48, bg_color=color, font_size=15)

    # Título principal (lado direito)
    txt(sl, "Contact Center", 4.5, 1.5, 8.5, 0.8,
        size=44, bold=True, color=TEXT_WHITE)
    txt(sl, "Data Lakehouse", 4.5, 2.25, 8.5, 0.8,
        size=44, bold=True, color=ACCENT_BLUE)
    txt(sl, "Pipeline end-to-end na AWS com arquitetura Medallion,\n"
            "modelagem Star Schema e orquestração com Apache Airflow.",
        4.5, 3.25, 8.5, 1.0, size=17, color=TEXT_LIGHT)

    box(sl, 4.5, 4.55, 8.3, 0.04, fill_color=ACCENT_TEAL)

    txt(sl, "Ilciley Junior  ·  Engenheiro de Dados",
        4.5, 4.75, 8.5, 0.5, size=15, color=TEXT_MUTED)
    txt(sl, "github.com/ilcileyjunior01",
        4.5, 5.2, 8.5, 0.4, size=13, color=ACCENT_BLUE)


def slide_02_agenda(prs):
    """Agenda / Sumário"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl)
    header(sl, "Agenda", "O que será apresentado")

    items = [
        ("01", "Problema de Negócio",        "Por que esse projeto existe?"),
        ("02", "Arquitetura de Dados",        "Stack AWS e camadas do Lakehouse"),
        ("03", "Data Quality",               "Tipagem, deduplicação e quarentena"),
        ("04", "Orquestração",               "Apache Airflow — DAGs e fluxo"),
        ("05", "Planos de Contingência",     "O que acontece quando um job falha"),
        ("06", "Resultados",                 "KPIs e dados reais carregados"),
        ("07", "Stack & Custo",              "Tecnologias e custo < $5/mês"),
    ]

    cols = [(0.5, 3.9), (6.8, 6.0)]
    rows_per_col = 4

    for i, (num, title, desc) in enumerate(items):
        col_idx = 0 if i < rows_per_col else 1
        row_idx = i if i < rows_per_col else i - rows_per_col
        x_base = cols[col_idx][0]
        w_card  = cols[col_idx][1]
        y = 1.45 + row_idx * 1.42

        box(sl, x_base, y, w_card, 1.18, fill_color=BG_CARD)
        txt(sl, num, x_base + 0.15, y + 0.12, 0.55, 0.5,
            size=22, bold=True, color=ACCENT_BLUE)
        txt(sl, title, x_base + 0.75, y + 0.1, w_card - 0.9, 0.38,
            size=15, bold=True, color=TEXT_WHITE)
        txt(sl, desc,  x_base + 0.75, y + 0.55, w_card - 0.9, 0.45,
            size=12, color=TEXT_MUTED)


def slide_03_problema(prs):
    """Problema de Negócio"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=ACCENT_ORG)
    header(sl, "Problema de Negócio",
           "Dados dispersos impossibilitam decisões ágeis no Contact Center")

    problems = [
        ("📦", "Dados Dispersos",
         "Chamadas, tickets, chats e WhatsApp\nem sistemas isolados. Nenhuma visão\nunificada do atendimento."),
        ("⚡", "Volume & Velocidade",
         "Milhares de eventos por dia. A\ninfraestrutura legada não acompanha\no crescimento da operação."),
        ("📊", "Decisões no Escuro",
         "Relatórios manuais em Excel,\ndemora de dias, erros frequentes\ne dados desatualizados."),
        ("🔒", "Conformidade LGPD",
         "Dados pessoais de clientes (CPF,\ne-mail, telefone) expostos sem\nmascaramento ou controle de acesso."),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        x = 0.5 + i * 3.22
        box(sl, x, 1.55, 3.0, 4.6, fill_color=BG_CARD)
        txt(sl, icon,  x + 0.15, 1.7,  2.7, 0.55, size=28, align=PP_ALIGN.CENTER)
        txt(sl, title, x + 0.15, 2.32, 2.7, 0.48,
            size=15, bold=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)
        box(sl, x + 0.3, 2.88, 2.4, 0.04, fill_color=ACCENT_ORG)
        txt(sl, desc,  x + 0.15, 3.0,  2.7, 1.9,
            size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Seta de solução
    txt(sl, "→  Solução: Data Lakehouse centralizado, automatizado e governado na AWS",
        0.5, 6.4, 12.3, 0.5, size=15, bold=True, color=ACCENT_TEAL)


def slide_04_arquitetura(prs):
    """Arquitetura de Dados — usa o diagrama PNG gerado pelo matplotlib"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=ACCENT_TEAL)
    header(sl, "Arquitetura de Dados",
           "Arquitetura Medallion: Fonte → Bronze → Silver → Gold → Analytics")

    # Prefere o diagrama com ícones AWS; cai no matplotlib se não existir
    base = os.path.dirname(__file__)
    diagram_path = os.path.join(base, "architecture_aws_icons.png")
    if not os.path.exists(diagram_path):
        diagram_path = os.path.join(base, "architecture_diagram.png")

    if os.path.exists(diagram_path):
        sl.shapes.add_picture(
            diagram_path,
            Inches(0.2), Inches(1.3),
            width=Inches(12.93), height=Inches(6.05)
        )
    else:
        txt(sl, "DIAGRAMA NAO ENCONTRADO — execute generate_aws_diagram.py primeiro",
            0.5, 3.5, 12.3, 0.6, size=14, color=RGBColor(0xFF, 0x00, 0x00),
            align=PP_ALIGN.CENTER)


def slide_05_data_quality(prs):
    """Data Quality"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=OK_GREEN)
    header(sl, "Data Quality",
           "Integridade garantida em cada camada do pipeline")

    practices = [
        ("Tipagem de Dados", OK_GREEN,
         ["Cast explícito em todos os jobs PySpark",
          "Timestamps: timestamp(6) → TimestampType()",
          "Flags booleanos: SMALLINT 0/1 (compatível Redshift)",
          "Numerais: nr_duracao → IntegerType / DoubleType",
          "Datas: DateType com format '%Y-%m-%d'"]),
        ("Deduplicação CDC", ACCENT_BLUE,
         ["Window.partitionBy(pk).orderBy(dt_cdc_evento DESC)",
          "row_number() == 1 → mantém evento mais recente",
          "Elimina duplicatas antes do MERGE",
          "Glue Job Bookmarks: impede releitura de arquivos",
          "Watermark JSON por tabela em S3/checkpoints/"]),
        ("Quarentena Automática", WARN_YELLOW,
         ["Registros com PK nula → isolados com motivo",
          "Timestamps ausentes → rejeitados com log",
          "Particionados por ano/mês/dia de ingestão",
          "Nunca causam falha no job principal",
          "Monitorados via CloudWatch Logs"]),
        ("MERGE ACID + Hash MD5", ACCENT_TEAL,
         ["hash_registro = MD5(todas colunas negócio)",
          "MERGE atualiza apenas se hash mudou",
          "Idempotente: reprocessar não duplica dados",
          "Iceberg v2: ACID, time travel, schema evolution",
          "Partição sobrescrita via overwritePartitions()"]),
    ]

    for i, (title, color, items) in enumerate(practices):
        x = 0.35 + i * 3.25
        box(sl, x, 1.5, 3.05, 5.35, fill_color=BG_CARD)
        box(sl, x, 1.5, 3.05, 0.45, fill_color=color)
        txt(sl, title, x + 0.1, 1.55, 2.85, 0.38,
            size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            txt(sl, f"• {item}", x + 0.12, 2.1 + j * 0.85, 2.82, 0.7,
                size=11, color=TEXT_LIGHT)


def slide_06_orquestracao(prs):
    """Orquestração"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=RGBColor(0x01, 0x7C, 0xEE))
    header(sl, "Orquestração — Apache Airflow 2.9.3",
           "2 DAGs orquestrando 40 Glue Jobs com paralelismo controlado")

    # DAG 1 — fluxo visual
    box(sl, 0.4, 1.5, 12.53, 3.25, fill_color=RGBColor(0x0A, 0x20, 0x35))
    txt(sl, "DAG 1: cc_pipeline_diario  [schedule: 0 2 * * *]",
        0.6, 1.55, 8.0, 0.4, size=13, bold=True, color=ACCENT_BLUE)
    txt(sl, "start_layer: bronze | silver | gold | redshift",
        8.7, 1.55, 4.0, 0.4, size=11, color=TEXT_MUTED)

    waves = [
        ("início",         "",                            BG_CARD,      0.55),
        ("Wave 1\nBronze→Silver\n18 jobs\n(paralelo)",   "", RGBColor(0x5C, 0x30, 0x00), 2.15),
        ("Wave 2\nDimensões Gold\n11 jobs\n(paralelo)",  "", RGBColor(0x1A, 0x4A, 0x1A), 4.2),
        ("Wave 3\nFatos Base\n7 jobs\n(paralelo)",       "", RGBColor(0x4A, 0x35, 0x00), 6.25),
        ("Wave 4\nFatos Dep.\n4 jobs\n(sequencial)",     "", RGBColor(0x3A, 0x00, 0x5C), 8.3),
        ("Trigger\ncc_carga_\nredshift",                 "", RGBColor(0x00, 0x35, 0x5C), 10.35),
    ]

    for label, _, color, x in waves:
        box(sl, x, 1.95, 1.7, 2.55, fill_color=color)
        txt(sl, label, x + 0.05, 2.05, 1.6, 2.35,
            size=11, bold=False, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        if x < 10.35:
            txt(sl, "→", x + 1.72, 2.95, 0.38, 0.55,
                size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    txt(sl, "ShortCircuitOperator (trigger_rule=ALL_DONE): gates permitem pular camadas sem quebrar o pipeline",
        0.6, 4.6, 12.1, 0.4, size=11, color=TEXT_MUTED)

    # DAG 2
    box(sl, 0.4, 5.1, 12.53, 1.75, fill_color=RGBColor(0x0A, 0x20, 0x35))
    txt(sl, "DAG 2: cc_carga_redshift  [schedule: None — acionada pela DAG 1 ou manualmente]",
        0.6, 5.18, 10.0, 0.4, size=13, bold=True, color=ACCENT_TEAL)

    dag2_steps = [
        "início_carga", "11 dims\n(paralelo)", "dims\ncarregadas",
        "11 fatos\n(paralelo)", "carga\nconcluída"
    ]
    for i, step in enumerate(dag2_steps):
        x2 = 0.6 + i * 2.45
        box(sl, x2, 5.65, 2.15, 1.0, fill_color=BG_CARD)
        txt(sl, step, x2 + 0.05, 5.72, 2.05, 0.85,
            size=11, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        if i < len(dag2_steps) - 1:
            txt(sl, "→", x2 + 2.17, 5.95, 0.25, 0.45,
                size=16, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    txt(sl, "Estratégia: Athena UNLOAD (Iceberg → CSV) → S3 staging → Redshift COPY  |  TRUNCATE + COPY por tabela",
        0.6, 6.75, 12.1, 0.38, size=11, color=TEXT_MUTED)


def slide_07_contingencia(prs):
    """Planos de Contingência"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=WARN_YELLOW)
    header(sl, "Planos de Contingência",
           "Resiliência end-to-end: detecção, isolamento e recuperação sem retrabalho")

    scenarios = [
        ("Job Bronze → Silver falha",
         WARN_YELLOW,
         ["CloudWatch Alarm dispara → SNS notifica",
          "Job Bookmark preserva posição: retry sem\nreprocessar o que já passou",
          "Trigger: start_layer=bronze para reprocessar\napenas aquela data",
          "Registros inválidos → Quarentena (não bloqueiam)"]),
        ("Job Silver → Gold falha",
         RGBColor(0xE7, 0x4C, 0x3C),
         ["ShortCircuitOperator: waves seguintes são\nmarcadas SKIPPED (não FAILED)",
          "Bronze intacto: dados brutos preservados",
          "Retry: start_layer=silver — reprocessa Gold\ne Redshift sem tocar Bronze",
          "MERGE idempotente: reexecutar não duplica"]),
        ("Carga Redshift falha",
         ACCENT_BLUE,
         ["DAG cc_carga_redshift pode ser acionada\nmanualmente a qualquer momento",
          "start_layer=redshift: pula todo o pipeline,\nexecuta só TRUNCATE + COPY",
          "Gold em S3 sempre disponível como fallback\n(Athena consulta direto se necessário)",
          "Auto-pause: workgroup não gera custo\nenquanto está em standby"]),
        ("Dados corrompidos na origem",
         RGBColor(0x8E, 0x44, 0xAD),
         ["Iceberg v2 Time Travel: SELECT AS OF\npermite voltar a versão anterior da tabela",
          "Bronze imutável: dado original nunca\né sobrescrito",
          "Watermark JSON: ajustar para reprocessar\na partir de qualquer data",
          "Quarentena auditável: rastreio de todos\nos registros rejeitados"]),
    ]

    for i, (title, color, items) in enumerate(scenarios):
        x = 0.35 + i * 3.25
        box(sl, x, 1.5, 3.05, 5.35, fill_color=BG_CARD)
        box(sl, x, 1.5, 3.05, 0.45, fill_color=color)
        txt(sl, title, x + 0.1, 1.55, 2.85, 0.38,
            size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            txt(sl, f"• {item}", x + 0.12, 2.1 + j * 1.18, 2.82, 1.0,
                size=11, color=TEXT_LIGHT)


def slide_08_modelo_dados(prs):
    """Modelo de Dados — Star Schema"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=ACCENT_ORG)
    header(sl, "Modelo de Dados — Star Schema Gold",
           "11 Dimensões + 11 Tabelas Fato — Iceberg v2, Parquet/Snappy, particionadas")

    # Dimensões
    box(sl, 0.35, 1.5, 5.9, 5.35, fill_color=BG_CARD)
    txt(sl, "DIMENSÕES (11)", 0.55, 1.58, 5.5, 0.38,
        size=14, bold=True, color=ACCENT_ORG)
    box(sl, 0.55, 2.02, 5.5, 0.04, fill_color=ACCENT_ORG)

    dims = [
        ("dim_data",              "sk_data = yyyyMMdd (INT)"),
        ("dim_cliente",           "4.159 clientes, PII mascarado"),
        ("dim_operador",          "208 operadores, sk_supervisor (auto-ref)"),
        ("dim_fila",              "17 filas, SLA em segundos/minutos"),
        ("dim_canal",             "Telefone · Chat · WhatsApp · Email"),
        ("dim_campanha",          "25 campanhas, vigência e flags"),
        ("dim_skill",             "26 skills, faixa de nível"),
        ("dim_status_chamada",    "fl_chamada_concluida, fl_abandonada"),
        ("dim_status_ticket",     "fl_aberto · fl_resolvido · fl_cancelado"),
        ("dim_categoria_ticket",  "5 categorias de suporte"),
        ("dim_prioridade_ticket", "CRITICA · ALTA · MEDIA · BAIXA"),
    ]
    for i, (name, desc) in enumerate(dims):
        y = 2.18 + i * 0.43
        txt(sl, name, 0.58, y, 2.55, 0.38,
            size=11, bold=True, color=ACCENT_ORG)
        txt(sl, desc, 3.15, y, 2.9, 0.38,
            size=10, color=TEXT_LIGHT)

    # Fatos
    box(sl, 6.65, 1.5, 6.3, 5.35, fill_color=BG_CARD)
    txt(sl, "TABELAS FATO (11)", 6.85, 1.58, 5.9, 0.38,
        size=14, bold=True, color=ACCENT_BLUE)
    box(sl, 6.85, 2.02, 5.9, 0.04, fill_color=ACCENT_BLUE)

    fatos = [
        ("fato_chamada",               "4.168 registros · sk_data_inicio (part.)"),
        ("fato_ticket",                "2.059 registros · SLA, TRT em minutos"),
        ("fato_qualidade",             "2.504 avaliações · nota, faixa, feedback"),
        ("fato_discagem",              "4.189 discagens · fl_atendida"),
        ("fato_ura_navegacao",         "4.126 navegações · fl_abandonou_ura"),
        ("fato_chat",                  "2.113 sessões · duração, fl_completo"),
        ("fato_whatsapp",              "1.386 atendimentos · nr_telefone_mask"),
        ("fato_jornada_operador",      "5.985 jornadas · horas, presença"),
        ("fato_metricas_operacionais", "4.187 métricas · TMA, TME, SLA fila"),
        ("fato_interacao_ticket",      "2.494 interações · canal, tamanho obs."),
        ("fato_mensagem_chat",         "4.905 mensagens · remetente, chars"),
    ]
    for i, (name, desc) in enumerate(fatos):
        y = 2.18 + i * 0.43
        txt(sl, name, 6.88, y, 3.1, 0.38,
            size=11, bold=True, color=ACCENT_BLUE)
        txt(sl, desc, 10.0, y, 2.75, 0.38,
            size=10, color=TEXT_LIGHT)

    # Rodapé
    txt(sl, "Registro sentinela sk = -1 (DESCONHECIDO) em todas as dimensões — integridade referencial garantida",
        0.35, 6.95, 12.6, 0.38, size=11, color=TEXT_MUTED)


def slide_09_resultados(prs):
    """Resultados"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=OK_GREEN)
    header(sl, "Resultados — Dados Reais Carregados",
           "Execução de referência: 2026-07-10 a 14  |  Pipeline 100% validado")

    # Métricas de destaque
    metrics = [
        ("40", "Jobs PySpark\nexecutados", ACCENT_BLUE),
        ("22", "Tabelas Gold\nno Redshift", ACCENT_TEAL),
        ("9",  "Views KPI\nvalidadas",     OK_GREEN),
        ("240","Testes\npassando",          ACCENT_ORG),
        ("0",  "Orphan\nrecords",           RGBColor(0x2E, 0xCC, 0x71)),
    ]
    for i, (num, label, color) in enumerate(metrics):
        x = 0.4 + i * 2.55
        box(sl, x, 1.5, 2.3, 1.55, fill_color=BG_CARD)
        txt(sl, num,   x + 0.1, 1.58, 2.1, 0.75,
            size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(sl, label, x + 0.1, 2.3,  2.1, 0.65,
            size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Views KPI
    box(sl, 0.4, 3.3, 12.53, 3.55, fill_color=BG_CARD)
    txt(sl, "Views KPI — Redshift Serverless (schema gold)",
        0.6, 3.38, 7.0, 0.38, size=13, bold=True, color=OK_GREEN)

    kpis = [
        ("vw_kpi_01", "Volume Chamadas",        "144 linhas",  "TMA abandono ~1 min, atendidas ~23 min"),
        ("vw_kpi_02", "Performance Operadores", "208 linhas",  "Top: 29 chamadas, TMA 18,36 min"),
        ("vw_kpi_03", "Qualidade",              "1.523 linhas","Nota média 7,0–10,0 por operador/mês"),
        ("vw_kpi_04", "Volume Tickets",         "1.350 linhas","Categorias: Elogio · Reclamação · Solicitação"),
        ("vw_kpi_05", "Eficiência Tickets",     "144 linhas",  "TRT médio 3,7–9,4h, % SLA cumprido"),
        ("vw_kpi_06", "Volume Digital",         "72 linhas",   "WhatsApp ~67 min vs Chat ~30 min"),
        ("vw_kpi_08", "Campanhas",              "25 linhas",   "Taxa de atendimento 55–66%"),
        ("vw_kpi_09", "Métricas de Fila",       "611 linhas",  "Nível de serviço 86–93% por fila"),
        ("vw_kpi_12", "Efetividade URA",        "252 linhas",  "Opção mais usada: 4 - FALAR COM OPERADOR"),
    ]

    # 3 colunas
    col_w = 4.1
    for i, (code, name, rows, insight) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.17
        y = 3.85 + row * 0.98
        txt(sl, f"{code}  ·  {name}", x, y, col_w - 0.15, 0.35,
            size=11, bold=True, color=TEXT_WHITE)
        txt(sl, f"{rows}  |  {insight}", x, y + 0.35, col_w - 0.15, 0.45,
            size=10, color=TEXT_MUTED)


def slide_10_lgpd(prs):
    """LGPD e Conformidade"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=RGBColor(0x8E, 0x44, 0xAD))
    header(sl, "Conformidade LGPD",
           "Dados pessoais mascarados de forma irreversível na camada Silver")

    # Tabela de mascaramento
    headers_row = ["Campo PII", "Técnica", "Exemplo", "Camada de Acesso"]
    rows_data = [
        ["CPF / Nr. Documento", "3 dígitos + ***** + 2 dígitos",          "123*****45",      "Bronze (restrito)"],
        ["E-mail",              "Domínio preservado, local mascarado",     "***@gmail.com",   "Bronze (restrito)"],
        ["Telefone",            "6 asteriscos + 4 últimos dígitos",        "******1234",      "Bronze (restrito)"],
        ["CEP",                 "5 primeiros + ***",                       "01310***",        "Bronze (restrito)"],
        ["Login do operador",   "3 primeiros chars + ***",                 "joh***",          "Bronze (restrito)"],
        ["Conteúdo de chat",    "Drop completo → nr_tamanho_chars",        "[removido]",      "Nunca salvo"],
        ["Obs. de ticket",      "Drop completo → nr_tamanho_obs_chars",    "[removido]",      "Nunca salvo"],
        ["Feedback avaliação",  "Drop completo → nr_tamanho_feedback_chars","[removido]",     "Nunca salvo"],
    ]

    col_widths = [2.8, 3.4, 2.0, 2.0]
    col_x      = [0.4, 3.22, 6.64, 8.66]
    header_y   = 1.52

    # cabeçalho da tabela
    box(sl, 0.4, header_y, 10.25, 0.42, fill_color=RGBColor(0x8E, 0x44, 0xAD))
    for j, h in enumerate(headers_row):
        txt(sl, h, col_x[j] + 0.08, header_y + 0.06, col_widths[j] - 0.1, 0.35,
            size=12, bold=True, color=TEXT_WHITE)

    for i, row in enumerate(rows_data):
        ry = header_y + 0.5 + i * 0.58
        bg_r = BG_CARD if i % 2 == 0 else RGBColor(0x10, 0x22, 0x38)
        box(sl, 0.4, ry, 10.25, 0.52, fill_color=bg_r)
        for j, cell in enumerate(row):
            color = TEXT_WHITE if j == 0 else TEXT_LIGHT
            if cell == "[removido]":
                color = RGBColor(0xE7, 0x4C, 0x3C)
            txt(sl, cell, col_x[j] + 0.08, ry + 0.08, col_widths[j] - 0.1, 0.38,
                size=11, color=color)

    # Controles
    box(sl, 10.8, 1.52, 2.13, 5.28, fill_color=BG_CARD)
    txt(sl, "Controles de\nAcesso", 10.9, 1.6, 1.95, 0.55,
        size=13, bold=True, color=RGBColor(0x8E, 0x44, 0xAD), align=PP_ALIGN.CENTER)

    controls = [
        "AWS Lake\nFormation",
        "Controle por\ncolunas",
        "Bronze: IAM\nrestrito",
        "CloudTrail:\nauditoria",
        "Silver+: PII\nmascarado",
    ]
    for i, c in enumerate(controls):
        box(sl, 10.88, 2.3 + i * 0.9, 1.95, 0.72, fill_color=RGBColor(0x3A, 0x10, 0x5C))
        txt(sl, c, 10.92, 2.36 + i * 0.9, 1.87, 0.6,
            size=10, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Dado original permanece APENAS no Bronze — acesso controlado via IAM + Lake Formation",
        0.4, 6.98, 12.53, 0.38, size=11, color=TEXT_MUTED)


def slide_11_stack_custo(prs):
    """Stack & Custo"""
    sl = blank_slide(prs)
    bg(sl)
    accent_bar(sl, color=ACCENT_ORG)
    header(sl, "Stack Tecnológica & Custo AWS",
           "Arquitetura serverless e pay-per-use — custo < $5/mês em escala de demonstração")

    # Stack
    stack_items = [
        ("Processamento",  "PySpark 3.3\nAWS Glue 4.0\nEMR Serverless",        ACCENT_ORG),
        ("Orquestração",   "Apache Airflow 2.9.3\nDocker Compose\nGlue Workflows", RGBColor(0x01, 0x7C, 0xEE)),
        ("Armazenamento",  "Amazon S3\nIceberg v2\nParquet + Snappy",            RGBColor(0x16, 0x9A, 0x45)),
        ("Catálogo",       "AWS Glue Data Catalog\nAmazon Athena\nLake Formation", RGBColor(0xFF, 0x66, 0x00)),
        ("Data Warehouse", "Redshift Serverless\n8 RPUs · auto-pause\nData API", RGBColor(0x8C, 0x4F, 0xFF)),
        ("IaC & CI/CD",    "Terraform\nGitHub Actions\n7 jobs CI pipeline",      RGBColor(0x62, 0x3C, 0xE4)),
    ]

    for i, (cat, tech, color) in enumerate(stack_items):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.17
        y = 1.5 + row * 2.2
        box(sl, x, y, 3.9, 1.95, fill_color=BG_CARD)
        box(sl, x, y, 3.9, 0.42, fill_color=color)
        txt(sl, cat, x + 0.12, y + 0.07, 3.65, 0.35,
            size=13, bold=True, color=TEXT_WHITE)
        txt(sl, tech, x + 0.12, y + 0.55, 3.65, 1.25,
            size=12, color=TEXT_LIGHT)

    # Custo
    box(sl, 0.4, 5.95, 12.53, 1.3, fill_color=BG_CARD)
    txt(sl, "Custo Estimado (escala de demonstração)",
        0.6, 6.02, 5.0, 0.38, size=13, bold=True, color=ACCENT_ORG)

    cost_items = [
        ("Ocioso (sem jobs)", "~$0,57/mês"),
        ("2 runs completos/mês", "~$5/mês"),
        ("Glue Jobs (40×G.1X)", "~$3,20"),
        ("Redshift auto-pause", "~$0,50"),
        ("S3 + Athena + outros", "~$0,13"),
    ]
    for i, (label, val) in enumerate(cost_items):
        x = 0.6 + i * 2.48
        txt(sl, label, x, 6.45, 2.3, 0.38, size=10, color=TEXT_MUTED)
        txt(sl, val,   x, 6.82, 2.3, 0.38, size=14, bold=True, color=ACCENT_ORG)


def slide_12_conclusao(prs):
    """Conclusão"""
    sl = blank_slide(prs)
    bg(sl, BG_DARK)

    box(sl, 0, 0, 4.2, 7.5, fill_color=RGBColor(0x09, 0x12, 0x1E))
    box(sl, 4.0, 0, 0.07, 7.5, fill_color=ACCENT_TEAL)

    txt(sl, "Conclusão", 4.5, 1.2, 8.5, 0.65,
        size=36, bold=True, color=TEXT_WHITE)
    box(sl, 4.5, 1.95, 8.3, 0.05, fill_color=ACCENT_TEAL)

    achievements = [
        ("✅", "Pipeline end-to-end executado",
               "40 jobs PySpark validados — Bronze → Silver → Gold → Redshift"),
        ("✅", "Data Quality garantido",
               "Deduplicação CDC, tipagem, quarentena, MERGE ACID idempotente"),
        ("✅", "Orquestração robusta",
               "Airflow 2.9.3 com reprocessamento por camada sem tocar a origem"),
        ("✅", "LGPD compliant",
               "PII mascarado irreversivelmente na Silver; Bronze controlado por IAM"),
        ("✅", "Resultados reais",
               "22 tabelas + 9 views KPI carregadas · 55K+ registros · 0 orphans"),
        ("✅", "Custo < $5/mês",
               "Serverless, auto-pause, pay-per-use em todas as camadas"),
    ]

    for i, (icon, title, desc) in enumerate(achievements):
        y = 2.15 + i * 0.82
        txt(sl, icon,  4.5,  y,       0.5,  0.65, size=18)
        txt(sl, title, 5.05, y,       7.7,  0.35, size=14, bold=True, color=ACCENT_TEAL)
        txt(sl, desc,  5.05, y + 0.36, 7.7, 0.35, size=12, color=TEXT_MUTED)

    txt(sl, "Ilciley Junior",
        4.5, 7.0, 5.0, 0.38, size=14, bold=True, color=TEXT_WHITE)
    txt(sl, "github.com/ilcileyjunior01  ·  linkedin.com/in/ilcileyjunior",
        4.5, 7.28, 8.5, 0.35, size=12, color=ACCENT_BLUE)

    # Lado esquerdo — destaques numéricos
    highlights = [
        ("40", "Glue Jobs"),
        ("22", "Tabelas Gold"),
        ("9",  "Views KPI"),
        ("240","Testes"),
    ]
    for i, (num, label) in enumerate(highlights):
        y = 1.5 + i * 1.38
        txt(sl, num,   0.4, y,       3.4, 0.72,
            size=38, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
        txt(sl, label, 0.4, y + 0.7, 3.4, 0.42,
            size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_problema(prs)
    slide_04_arquitetura(prs)
    slide_05_data_quality(prs)
    slide_06_orquestracao(prs)
    slide_07_contingencia(prs)
    slide_08_modelo_dados(prs)
    slide_09_resultados(prs)
    slide_10_lgpd(prs)
    slide_11_stack_custo(prs)
    slide_12_conclusao(prs)

    out = "docs/contact_center_data_lakehouse.pptx"
    prs.save(out)
    print(f"OK  Apresentacao salva em: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
